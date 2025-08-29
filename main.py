import os
import io
import traceback
import google.generativeai as genai
import docx # Keep using python-docx
from docx.document import Document as DocxDocument # To handle type hinting
import PyPDF2
from PIL import Image
from pptx import Presentation
from flask import Flask, request, jsonify, send_file
from flask_cors import CORS
from werkzeug.utils import secure_filename

# --- Settings ---
app = Flask(__name__, static_folder='.', static_url_path='')
CORS(app)

# --- Gemini API Setup ---
model = None
api_key_error = None
try:
    GEMINI_API_KEY = os.environ.get('GEMINI_API_KEY')
    if not GEMINI_API_KEY:
        raise ValueError("GEMINI_API_KEY environment variable not set or found.")
    
    genai.configure(api_key=GEMINI_API_KEY)
    model = genai.GenerativeModel('gemini-1.5-pro-latest') 
    print("✅ Gemini API configured successfully with gemini-1.5-pro-latest model.")
except Exception as e:
    api_key_error = str(e)
    print(f"!!!!!! FATAL ERROR during Gemini API setup: {api_key_error}")

# --- Helper function for making translation calls to the API ---
def translate_text_api(text_to_translate, target_lang):
    """Sends a single text chunk to Gemini for translation and returns the result."""
    if not text_to_translate or not text_to_translate.strip():
        return ""
    try:
        # Simple prompt for individual text chunks
        prompt = f"Translate the following text to {target_lang}. Provide only the translated text, nothing else:\n\n{text_to_translate}"
        response = model.generate_content(prompt)
        # Using .strip() to remove potential leading/trailing whitespace from the API response
        return response.text.strip()
    except Exception as e:
        print(f"--- API translation error for a chunk: {e}")
        # Return original text if translation fails for this specific chunk
        return text_to_translate

# --- NEW ADVANCED DOCX TRANSLATION FUNCTION ---
def translate_docx_in_place(doc: DocxDocument, target_lang: str):
    """
    Translates a docx document in-place, preserving formatting.
    It iterates through paragraphs and table cells, translates their text,
    and replaces the original text.
    """
    print("Starting in-place DOCX translation...")
    
    # 1. Translate Paragraphs
    for para in doc.paragraphs:
        if para.text.strip(): # Check if there is text to translate
            original_text = para.text
            print(f"  Translating paragraph: '{original_text[:50]}...'")
            translated_text = translate_text_api(original_text, target_lang)
            
            # Replace the paragraph's text while preserving its style
            para.text = translated_text

    # 2. Translate Tables
    print("Scanning for tables...")
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                # Each cell can contain multiple paragraphs
                for para in cell.paragraphs:
                    if para.text.strip():
                        original_text = para.text
                        print(f"  Translating table cell: '{original_text[:50]}...'")
                        translated_text = translate_text_api(original_text, target_lang)
                        para.text = translated_text
                        
    print("In-place DOCX translation finished.")
    return doc

# --- Text Extraction for other file types (unchanged) ---
def read_text_from_pdf(stream):
    reader = PyPDF2.PdfReader(stream)
    return '\n'.join([page.extract_text() for page in reader.pages if page.extract_text()])

def read_text_from_pptx(stream):
    prs = Presentation(stream)
    text_runs = []
    # ... (rest of the function is the same)
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        text_runs.append(run.text)
    return '\n'.join(text_runs)


def create_docx_from_text(text):
    # This function is now only a fallback for PDF/PPTX
    mem_file = io.BytesIO()
    doc = docx.Document()
    doc.add_paragraph(text)
    doc.save(mem_file)
    mem_file.seek(0)
    return mem_file

# --- Main Application Routes ---

@app.route('/')
def serve_index():
    return app.send_static_file('index.html')

@app.route('/translate-file', methods=['POST'])
def translate_file_handler():
    print("--- Received request for /translate-file ---")
    if not model:
        return jsonify({"error": "API service is not configured."}), 500
    if 'file' not in request.files:
        return jsonify({"error": "No file part in the request."}), 400

    file = request.files['file']
    if file.filename == '':
        return jsonify({"error": "No file selected."}), 400
    
    filename = secure_filename(file.filename).lower()
    print(f"Processing file: {filename}")

    try:
        target_lang = request.form.get('target_lang', 'English')
        new_filename = f"translated_{os.path.splitext(file.filename)[0]}.docx"
        
        # --- LOGIC BRANCHING BASED ON FILE TYPE ---

        if filename.endswith('.docx'):
            # Use the new advanced method for DOCX files
            original_doc = docx.Document(file.stream)
            translated_doc = translate_docx_in_place(original_doc, target_lang)
            
            # Save the modified document to a memory stream
            mem_file = io.BytesIO()
            translated_doc.save(mem_file)
            mem_file.seek(0)
            
            return send_file(
                mem_file,
                as_attachment=True,
                download_name=new_filename,
                mimetype='application/vnd.openxmlformats-officedocument.wordprocessingml.document'
            )

        # For all other file types, use the old method (text extraction)
        translated_text = ""
        
        if filename.endswith(('.png', '.jpg', '.jpeg')):
            image = Image.open(file.stream)
            prompt_parts = ["Extract and translate the text in this image to {target_lang}. Provide only the translation.", image]
            response = model.generate_content(prompt_parts)
            translated_text = response.text

        elif filename.endswith('.pdf'):
            original_text = read_text_from_pdf(file.stream)
            translated_text = translate_text_api(original_text, target_lang)

        elif filename.endswith('.pptx'):
            original_text = read_text_from_pptx(file.stream)
            translated_text = translate_text_api(original_text, target_lang)
        
        else:
            return jsonify({"error": "Unsupported file type."}), 400

        # Create a new docx for these non-docx file types
        translated_doc_stream = create_docx_from_text(translated_text)
        return send_file(
            translated_doc_stream,
            as_attachment=True,
            download_name=new_filename,
            mimetype='application/vnd.openxmlformats-officedocument.wordprocessingml.document'
        )

    except Exception as e:
        print(f"!!!!!! CRITICAL ERROR during file processing for {filename}: {e}")
        traceback.print_exc()
        return jsonify({"error": "An internal server error occurred during file processing."}), 500


@app.route('/translate-text', methods=['POST'])
def translate_text_handler():
    # This function is working correctly, no changes needed.
    # ... (code for this function is unchanged) ...
    if not model: return jsonify({"error": "API service is not configured."}), 500
    data = request.get_json()
    text = data.get('text')
    if not text: return jsonify({"error": "No text provided."}), 400
    target_lang = data.get('target_lang', 'English')
    try:
        prompt = (f"Translate the following text to '{target_lang}'. Provide only the professional translated text.\n\n{text}")
        response = model.generate_content(prompt)
        return jsonify({"translated_text": response.text})
    except Exception as e:
        print(f"!!!!!! An error occurred during text translation: {e}")
        return jsonify({"error": "An internal error occurred during translation."}), 500

if __name__ == "__main__":
    port = int(os.environ.get('PORT', 8080))
    app.run(debug=False, host='0.0.0.0', port=port)
